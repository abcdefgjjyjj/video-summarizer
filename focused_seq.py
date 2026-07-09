#!/usr/bin/env python3
"""
Sequential batch — reliable single-process transcription + synthesis.
Targets only relevant folders: Motion/Vision/I/O/Mapping/EAP/软件框架/半导体工艺.
"""
import os, sys, time, traceback
from pathlib import Path

PROJECT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(PROJECT_DIR))

from src.audio import extract_audio, get_audio_duration
from src.transcribe import transcribe, format_transcript_with_timestamps

API_KEY = "YOUR_API_KEY"
MODEL_SIZE = "base"
LLM_MODEL = "llm-chat"
HF_MIRROR = True
LANGUAGE = "zh"

RELEVANT_PATTERNS = [
    "ACS应用介绍(二)", "电机特殊运动", "1210时序规划与分析", "FC1250时序分析",
    "相机镜头建模", "新版视觉2.0", "AI技术在芯片检测", "视觉2新功能",
    "视觉基础知识", "MIT视觉GP", "视觉GP软件培训-员工H",
    "FC1250所用传感器", "凌华卡与雷赛卡", "基于1210的通讯",
    "建模注意事项",
    "EAP介绍和调试", "EAP的软件安装", "EAP模拟器",
    "设备稳定性能指标参数", "svn提交规则与1250贴片时序",
    "设备模块的封装和应用", "1220框架讲解",
    "Pick&Place工艺讲解(一)", "Pick&Place工艺讲解(二)", "Pick&Place工艺讲解(三)",
    "BGA封装流程", "Memory工艺",
    "1250设备常见问题及其分析方法（一）",
    "1250设备常见问题及其分析方法（二）",
    "1250设备常见问题及其分析方法（三）",
]

TOPIC_PROMPT = """请根据以下培训视频转录{pptx_part}，编写一份聚焦于以下技术领域的培训学习文档。

## 关注领域
本总结只关注以下内容，忽略与这些领域无关的部分：
- Motion（运动控制）：轴运动、运动轨迹规划、电机控制、时序分析、ACS运动控制器等
- Vision（视觉）：视觉硬件、相机镜头、图像处理、视觉标定、AI检测、GP软件等
- I/O（输入输出）：传感器、通讯协议、数据采集卡（凌华/雷赛）、电气IO等
- Mapping（映射/标定）：坐标系标定、建模、空间变换等
- EAP（设备自动化程序）：EAP架构、SECS/GEM通讯、EAP调试、EAP模拟器等
- EFEM（设备前端模块）：晶圆传输、前端模块控制等
- 用户管理：用户权限、角色管理、登录认证等
- 日志管理：日志记录、日志查看、日志分析等
- 参数文件管理：参数配置、配方管理、文件格式、版本控制等
- 半导体工艺知识：封装测试流程、Die Bond / Wire Bond / Flip Chip工艺、BGA/Memory封装等

## 输出结构
1. **培训概述**：主题、设备型号、培训目标（仅限上述关注领域）
2. **Motion / 运动控制**：轴配置、运动模式、时序规划、特殊运动等
3. **Vision / 视觉系统**：硬件配置、相机标定、图像处理流程、检测算法等
4. **I/O / 输入输出**：传感器类型、通讯方式、数据采集等
5. **EAP / Mapping / EFEM**：相关配置和流程
6. **软件框架**：用户管理、日志管理、参数文件管理、模块封装等
7. **半导体工艺**：封装流程、工艺参数、关键控制点等
8. **关键参数汇总表**
9. **关键结论/FAQ**

## 注意事项
- 禁止出现"作为...专家"、"好的，我将..."等AI味开场白，直接输出文档内容
- 只提取与上述关注领域相关的内容，无关部分直接跳过
- 如果某个关注领域在培训中没有涉及，该章节可以省略
- 不要忽略任何细节和参数说明
- 结合PPT中的术语和框架（如有）
- 不要编造，纯从培训材料中提取
- 中文输出，专业术语保留英文
- 如培训材料完全无相关内容，输出"本培训不涉及目标领域的内容"

---

## 视频转录文本

{transcript}

{pptx_content}
"""


def find_file(folder, *extensions):
    for f in sorted(os.listdir(folder)):
        if any(f.lower().endswith(ext) for ext in extensions):
            return os.path.join(folder, f)
    return None


def extract_pptx_text(pptx_path):
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
        return "\n".join(lines)[:8000]
    except Exception as e:
        return f"[PPTX extraction error: {e}]"


def synthesize(video_name, transcript, ppt_text, output_path):
    import httpx

    trans_snippet = transcript[:25000] if len(transcript) > 25000 else transcript
    pptx_part = "、PPT讲义" if ppt_text else ""
    pptx_content = f"\n\n---\n## PPT讲义内容\n\n{ppt_text}\n" if ppt_text else ""

    prompt = TOPIC_PROMPT.format(
        pptx_part=pptx_part,
        transcript=trans_snippet,
        pptx_content=pptx_content,
    )

    print(f"       Prompt: {len(prompt):,} chars, calling LLM...")

    r = httpx.post(
        "https://api.llm.com/v1/chat/completions",
        headers={"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"},
        json={"model": LLM_MODEL, "max_tokens": 8192, "messages": [{"role": "user", "content": prompt}]},
        timeout=httpx.Timeout(600.0, connect=30.0, read=300.0, write=30.0),
    )

    if r.status_code != 200:
        print(f"       [ERROR] API {r.status_code}: {r.text[:200]}")
        return None

    result = r.json()["choices"][0]["message"]["content"]
    Path(output_path).write_text(result, encoding="utf-8")
    return result


def process_folder(folder_path):
    folder_name = os.path.basename(folder_path)
    print(f"\n{'='*60}")
    print(f"[{folder_name}]")
    print(f"{'='*60}")

    # Check if already done
    existing_doc = find_file(folder_path, "_培训文档.md", "培训文档.md")
    if existing_doc:
        print(f"  [SKIP] Document exists: {os.path.basename(existing_doc)}")
        return True

    video_file = find_file(folder_path, ".mp4", ".mkv", ".mov", ".avi")
    pptx_file = find_file(folder_path, ".pptx", ".ppt")

    if not video_file:
        print(f"  [WARN] No video, checking for PPT-only...")
        if pptx_file:
            ppt_text = extract_pptx_text(pptx_file)
            print(f"  PPT: {os.path.basename(pptx_file)} ({len(ppt_text) if ppt_text else 0} chars)")
            stem = Path(pptx_file).stem
            doc_output = os.path.join(folder_path, f"{stem}_培训文档.md")
            result = synthesize(folder_name, "(无视频转录)", ppt_text, doc_output)
            if result:
                print(f"  [DONE] PPT-only doc: {len(result):,} chars")
                return True
            return False
        else:
            print(f"  [SKIP] No video, no PPT")
            return True

    print(f"  Video: {os.path.basename(video_file)} ({os.path.getsize(video_file)/1024/1024:.1f}MB)")
    if pptx_file:
        print(f"  PPT:   {os.path.basename(pptx_file)}")

    stem = Path(video_file).stem

    # Check for existing transcript
    transcript_file = find_file(folder_path, "_transcript.txt")
    transcript_text = None

    if transcript_file and not transcript_file.endswith("_timed.txt"):
        transcript_text = Path(transcript_file).read_text(encoding="utf-8")
        print(f"  [Reuse] Existing transcript: {len(transcript_text):,} chars")
    else:
        # Step 1: Extract audio
        print(f"  [1/3] Extracting audio...")
        t0 = time.time()
        audio_path = extract_audio(video_file, output_dir=folder_path)
        duration = get_audio_duration(audio_path)
        print(f"       {duration:.0f}s ({duration/60:.1f}min) in {time.time()-t0:.0f}s")

        # Step 2: Transcribe
        print(f"  [2/3] Transcribing (model={MODEL_SIZE})...")
        t0 = time.time()
        transcript = transcribe(audio_path, model_size=MODEL_SIZE, language=LANGUAGE, hf_mirror=HF_MIRROR)
        transcript_file = os.path.join(folder_path, f"{stem}_transcript.txt")
        timed_file = os.path.join(folder_path, f"{stem}_transcript_timed.txt")
        Path(transcript_file).write_text(transcript["text"], encoding="utf-8")
        Path(timed_file).write_text(format_transcript_with_timestamps(transcript["segments"]), encoding="utf-8")
        transcript_text = transcript["text"]
        print(f"       {len(transcript_text):,} chars, {len(transcript['segments'])} segments, {time.time()-t0:.0f}s")

        # Cleanup audio
        if os.path.exists(audio_path):
            os.remove(audio_path)

    # Extract PPTX if available
    ppt_text = ""
    if pptx_file:
        ppt_text = extract_pptx_text(pptx_file)
        print(f"  PPT:   {len(ppt_text)} chars extracted")

    # Step 3: Synthesize
    print(f"  [3/3] Synthesizing...")
    doc_output = os.path.join(folder_path, f"{stem}_培训文档.md")
    result = synthesize(folder_name, transcript_text, ppt_text, doc_output)

    if result:
        print(f"  [DONE] {len(result):,} chars")
        return True
    else:
        print(f"  [FAIL] Synthesis failed")
        return False


def main():
    base = sys.argv[1] if len(sys.argv) > 1 else r"<training-videos-directory>"
    total_start = time.time()

    dirs = sorted([
        os.path.join(base, d) for d in os.listdir(base)
        if os.path.isdir(os.path.join(base, d))
    ])

    # Filter relevant
    def is_relevant(name):
        for pat in RELEVANT_PATTERNS:
            if pat in name:
                return True
        return False

    relevant = []
    skipped_done = []
    skipped_irrelevant = []

    for folder_path in dirs:
        name = os.path.basename(folder_path)
        if not is_relevant(name):
            skipped_irrelevant.append(name)
            continue
        existing = find_file(folder_path, "_培训文档.md", "培训文档.md")
        if existing:
            skipped_done.append(name)
            continue
        relevant.append(folder_path)

    print(f"{'='*60}")
    print(f"Focused Sequential Batch Processing")
    print(f"{'='*60}")
    print(f"Total folders in dir: {len(dirs)}")
    print(f"Relevant & unprocessed: {len(relevant)}")
    print(f"Already done: {len(skipped_done)}")
    print(f"Irrelevant (skipped): {len(skipped_irrelevant)}")
    print(f"Model: {MODEL_SIZE} | LLM: {LLM_MODEL}")
    print(f"{'='*60}")

    if not relevant:
        print("\nAll done!")
        return

    success = 0
    failed = []
    for i, folder_path in enumerate(relevant):
        print(f"\n[{i+1}/{len(relevant)}]", end="")
        try:
            ok = process_folder(folder_path)
            if ok:
                success += 1
            else:
                failed.append(os.path.basename(folder_path))
        except Exception as e:
            print(f"\n  [FATAL] {e}")
            traceback.print_exc()
            failed.append(os.path.basename(folder_path))

    elapsed = time.time() - total_start
    print(f"\n{'='*60}")
    print(f"COMPLETE: {success}/{len(relevant)} done in {elapsed/60:.1f}min")
    if failed:
        print(f"Failed ({len(failed)}):")
        for f in failed:
            print(f"  - {f}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
