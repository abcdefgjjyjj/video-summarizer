#!/usr/bin/env python3
"""
Focused batch processing — only videos related to: Motion, Vision, I/O, Mapping,
EAP, EFEM, 用户管理, 日志管理, 参数文件管理, 半导体工艺知识.

Phase 1: 2-process parallel transcription
Phase 2: 6-thread parallel synthesis with topic-filtered prompt
"""
import os, sys, time
from pathlib import Path
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor, as_completed

PROJECT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(PROJECT_DIR))

API_KEY = "YOUR_API_KEY"
MODEL_SIZE = "base"
LLM_MODEL = "llm-chat"
HF_MIRROR = True
LANGUAGE = "zh"

# ---- Relevant folder keywords ----
RELEVANT_PATTERNS = [
    "ACS应用介绍(二)",           # Motion
    "电机特殊运动",              # Motion
    "1210时序规划与分析",        # Motion
    "FC1250时序分析",            # Motion
    "相机镜头建模",              # Vision
    "新版视觉2.0",               # Vision
    "AI技术在芯片检测",          # Vision
    "视觉2新功能",               # Vision + Mapping
    "视觉基础知识",              # Vision
    "MIT视觉GP",                 # Vision
    "视觉GP软件培训-员工H",      # Vision
    "FC1250所用传感器",          # I/O
    "凌华卡与雷赛卡",            # I/O
    "基于1210的通讯",            # I/O
    "建模注意事项",              # Mapping
    "EAP介绍和调试",             # EAP
    "EAP的软件安装",             # EAP
    "EAP模拟器",                 # EAP
    "设备稳定性能指标参数",       # 参数文件管理
    "svn提交规则与1250贴片时序", # 参数文件管理
    "设备模块的封装和应用",       # 软件框架
    "1220框架讲解",              # 软件框架
    "Pick&Place工艺讲解(一)",    # 半导体工艺
    "Pick&Place工艺讲解(二)",    # 半导体工艺
    "Pick&Place工艺讲解(三)",    # 半导体工艺
    "BGA封装流程",               # 半导体工艺
    "Memory工艺",                # 半导体工艺
    "1250设备常见问题及其分析方法（一）",  # 综合
    "1250设备常见问题及其分析方法（二）",  # 综合
    "1250设备常见问题及其分析方法（三）",  # 综合
]


def is_relevant(folder_name: str) -> bool:
    for pat in RELEVANT_PATTERNS:
        if pat in folder_name:
            return True
    return False


def _find_video(folder_path: str) -> str | None:
    for f in sorted(os.listdir(folder_path)):
        if any(f.lower().endswith(ext) for ext in ('.mp4', '.mkv', '.mov', '.avi')):
            return os.path.join(folder_path, f)
    return None


# ---- Phase 1 worker: extract audio + transcribe (CPU-bound) ----

def transcribe_worker(folder_path: str) -> dict:
    import sys
    from pathlib import Path
    sys.path.insert(0, str(Path(folder_path).parent.parent / "video-summarizer"))
    from src.audio import extract_audio, get_audio_duration
    from src.transcribe import transcribe, format_transcript_with_timestamps

    folder_name = os.path.basename(folder_path)
    video_file = _find_video(folder_path)

    stem = None
    if video_file:
        stem = Path(video_file).stem
    else:
        # No video — just mark for PPTX-only processing
        return {"folder": folder_name, "status": "no_video", "stem": folder_name}

    doc_file = os.path.join(folder_path, f"{stem}_培训文档.md")
    if os.path.exists(doc_file):
        return {"folder": folder_name, "status": "skipped", "reason": "doc exists"}

    t0 = time.time()
    try:
        audio_path = extract_audio(video_file, output_dir=folder_path)
        duration = get_audio_duration(audio_path)
        transcript = transcribe(audio_path, model_size=MODEL_SIZE, language=LANGUAGE, hf_mirror=HF_MIRROR)

        transcript_file = os.path.join(folder_path, f"{stem}_transcript.txt")
        timed_file = os.path.join(folder_path, f"{stem}_transcript_timed.txt")
        Path(transcript_file).write_text(transcript["text"], encoding="utf-8")
        Path(timed_file).write_text(format_transcript_with_timestamps(transcript["segments"]), encoding="utf-8")

        if os.path.exists(audio_path):
            os.remove(audio_path)

        elapsed = time.time() - t0
        return {
            "folder": folder_name, "status": "transcribed",
            "chars": len(transcript["text"]),
            "duration_min": round(duration / 60, 1),
            "elapsed_s": round(elapsed, 0), "stem": stem,
        }
    except Exception as e:
        return {"folder": folder_name, "error": str(e)}


# ---- Phase 2 worker: synthesize with topic-filtered prompt (I/O-bound) ----

TOPIC_FILTER_PROMPT = """请根据以下培训视频转录{f'}、PPT讲义' if ppt_text else ''}，编写一份聚焦于以下技术领域的培训学习文档。

## 关注领域
本总结只关注以下内容，**忽略与这些领域无关的部分**：
- **Motion（运动控制）**：轴运动、运动轨迹规划、电机控制、时序分析、ACS运动控制器等
- **Vision（视觉）**：视觉硬件、相机镜头、图像处理、视觉标定、AI检测、GP软件等
- **I/O（输入输出）**：传感器、通讯协议、数据采集卡（凌华/雷赛）、电气IO等
- **Mapping（映射/标定）**：坐标系标定、建模、空间变换等
- **EAP（设备自动化程序）**：EAP架构、SECS/GEM通讯、EAP调试、EAP模拟器等
- **EFEM（设备前端模块）**：晶圆传输、前端模块控制等
- **用户管理**：用户权限、角色管理、登录认证等
- **日志管理**：日志记录、日志查看、日志分析等
- **参数文件管理**：参数配置、配方管理、文件格式、版本控制等
- **半导体工艺知识**：封装测试流程、Die Bond / Wire Bond / Flip Chip工艺、BGA/Memory封装等

## 输出结构
1. **培训概述**：主题、设备型号、培训目标（仅限上述关注领域）
2. **Motion / 运动控制**：轴配置、运动模式、时序规划、特殊运动等
3. **Vision / 视觉系统**：硬件配置、相机标定、图像处理流程、检测算法等
4. **I/O / 输入输出**：传感器类型、通讯方式、数据采集等
5. **EAP / Mapping / EFEM**：相关配置和流程
6. **软件框架**：用户管理、日志管理、参数文件管理、模块封装等
7. **半导体工艺**：封装流程、工艺参数、关键控制点等
8. **关键参数汇总表**
9. **关键结论/FAQ**

## 注意事项
- 禁止出现"作为...专家"、"好的，我将..."等AI味开场白，直接输出文档内容
- 只提取与上述关注领域相关的内容，无关部分直接跳过
- 如果某个关注领域在培训中没有涉及，该章节可以省略
- 不要忽略任何细节和参数说明
- 结合PPT中的术语和框架（如有）
- 不要编造，纯从培训材料中提取
- 中文输出，专业术语保留英文
- 培训材料中完全没有任何相关内容的章节，输出"本培训不涉及目标领域的内容"

---

## 视频转录文本

{trans_snippet}

{f'---## PPT讲义内容{ppt_text}---' if ppt_text else ''}
"""


def synthesize_worker(args: tuple) -> dict:
    folder_path, api_key, model = args
    import httpx

    folder_name = os.path.basename(folder_path)

    # Extract PPTX/DOCX text
    ppt_text = ""
    for f in sorted(os.listdir(folder_path)):
        fp = os.path.join(folder_path, f)
        if f.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(fp)
                lines = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    lines.append(t)
                ppt_text = "\n".join(lines)[:8000]
            except Exception:
                pass
        elif f.endswith(".docx"):
            try:
                from docx import Document
                doc = Document(fp)
                parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        parts.append(para.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        parts.append(" | ".join(cell.text for cell in row.cells))
                ppt_text += "\n\n## 试题\n" + "\n".join(parts)
            except Exception:
                pass

    # Find transcript
    transcript_file = None
    stem = None
    for f in sorted(os.listdir(folder_path)):
        if f.endswith("_transcript.txt") and not f.endswith("_timed.txt"):
            transcript_file = os.path.join(folder_path, f)
            stem = Path(f).stem.replace("_transcript", "")
            break

    if not transcript_file:
        # No transcript — video-only? try PPT-only
        if ppt_text:
            stem = folder_name
            trans_snippet = "(无视频转录，仅PPT讲义内容)"
        else:
            return {"folder": folder_name, "error": "No transcript, no PPT"}
    else:
        transcript_text = Path(transcript_file).read_text(encoding="utf-8")
        trans_snippet = transcript_text[:25000] if len(transcript_text) > 25000 else transcript_text

    prompt = TOPIC_FILTER_PROMPT.format(
        trans_snippet=trans_snippet,
        ppt_text=ppt_text,
    )

    try:
        r = httpx.post(
            "https://api.llm.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": model, "max_tokens": 8192, "messages": [{"role": "user", "content": prompt}]},
            timeout=httpx.Timeout(600.0, connect=30.0, read=300.0, write=30.0),
        )
        if r.status_code != 200:
            return {"folder": folder_name, "error": f"API {r.status_code}: {r.text[:100]}"}

        result = r.json()["choices"][0]["message"]["content"]
        doc_output = os.path.join(folder_path, f"{stem}_培训文档.md")
        Path(doc_output).write_text(result, encoding="utf-8")
        return {"folder": folder_name, "status": "done", "chars": len(result)}
    except Exception as e:
        return {"folder": folder_name, "error": str(e)}


# ---- main ----

def main():
    base = sys.argv[1] if len(sys.argv) > 1 else r"<training-videos-directory>"

    dirs = sorted([
        os.path.join(base, d) for d in os.listdir(base)
        if os.path.isdir(os.path.join(base, d))
    ])

    # Filter: relevant + not yet processed
    unprocessed = []
    skipped_done = []
    skipped_irrelevant = []

    for folder_path in dirs:
        folder_name = os.path.basename(folder_path)

        if not is_relevant(folder_name):
            skipped_irrelevant.append(folder_name)
            continue

        has_doc = any(
            f.endswith("_培训文档.md") or f.endswith("培训文档.md")
            for f in os.listdir(folder_path)
        )
        if has_doc:
            skipped_done.append(folder_name)
            continue

        unprocessed.append(folder_path)

    print(f"{'='*60}")
    print(f"Focused Batch — Motion/Vision/I/O/Mapping/EAP/软件框架/工艺")
    print(f"{'='*60}")
    print(f"Total folders: {len(dirs)}")
    print(f"Relevant & unprocessed: {len(unprocessed)}")
    print(f"Already done: {len(skipped_done)}")
    print(f"Irrelevant (skipped): {len(skipped_irrelevant)}")

    if not unprocessed:
        print("\nAll relevant folders already processed!")
        return

    print(f"\nUnprocessed relevant folders:")
    for f in unprocessed:
        name = os.path.basename(f)
        has_vid = "VID" if _find_video(f) else "PPT-only"
        print(f"  - {name} [{has_vid}]")

    # ---- Phase 1: Parallel transcription (2 workers) ----
    print(f"\n{'='*60}")
    print("Phase 1: Parallel transcription (2 workers)")
    print(f"{'='*60}")
    t0 = time.time()

    with ProcessPoolExecutor(max_workers=2) as executor:
        futures = {executor.submit(transcribe_worker, f): f for f in unprocessed}
        for future in as_completed(futures):
            r = future.result()
            if r.get("status") == "transcribed":
                print(f"  [{r['folder']}] OK  {r['duration_min']}min, {r['chars']:,} chars, {r['elapsed_s']:.0f}s")
            elif r.get("status") == "skipped":
                print(f"  [{r['folder']}] SKIP ({r.get('reason')})")
            elif r.get("status") == "no_video":
                print(f"  [{r['folder']}] NO_VIDEO (will synthesize from PPT only)")
            else:
                print(f"  [{r['folder']}] FAIL: {r.get('error', 'unknown')}")

    t1 = time.time()
    print(f"Phase 1 done: {t1 - t0:.0f}s")

    # ---- Phase 2: Parallel synthesis (6 threads) ----
    print(f"\n{'='*60}")
    print("Phase 2: Parallel synthesis (6 threads)")
    print(f"{'='*60}")

    with ThreadPoolExecutor(max_workers=6) as executor:
        futures = {
            executor.submit(synthesize_worker, (f, API_KEY, LLM_MODEL)): f
            for f in unprocessed
        }
        for future in as_completed(futures):
            r = future.result()
            if r.get("status") == "done":
                print(f"  [{r['folder']}] DONE  {r['chars']:,} chars")
            else:
                print(f"  [{r['folder']}] FAIL: {r.get('error', 'unknown')}")

    t2 = time.time()
    print(f"\nPhase 2 done: {t2 - t1:.0f}s")
    print(f"Total: {t2 - t0:.0f}s")


if __name__ == "__main__":
    main()
