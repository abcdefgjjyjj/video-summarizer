#!/usr/bin/env python3
"""
批量培训视频处理管线。
用法: python batch_process.py "<training-videos-directory>"
"""

import os
import sys
import json
import time
import subprocess
from pathlib import Path

# 加入项目路径
sys.path.insert(0, str(Path(__file__).resolve().parent))

from src.audio import extract_audio, get_audio_duration
from src.transcribe import transcribe, format_transcript_with_timestamps
from src.llm import summarize

# ---------- config ----------
API_KEY = "YOUR_API_KEY"
MODEL_SIZE = "base"
LLM_MODEL = "llm-chat"
HF_MIRROR = True
LANGUAGE = "zh"
# ---------------------------

def find_file(folder, *extensions):
    """在文件夹中找第一个匹配扩展名的文件。"""
    for f in sorted(os.listdir(folder)):
        if any(f.lower().endswith(ext) for ext in extensions):
            return os.path.join(folder, f)
    return None

def extract_pptx(pptx_path, output_path):
    """提取 PPTX 文本。"""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
        text = "\n".join(lines)
        Path(output_path).write_text(text, encoding="utf-8")
        return text
    except Exception as e:
        print(f"  [WARN] PPTX extraction failed: {e}")
        return ""

def extract_docx(docx_path, output_path):
    """提取 DOCX 文本（试题）。"""
    try:
        from docx import Document
        doc = Document(docx_path)
        lines = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                lines.append(t)
        for table in doc.tables:
            for row in table.rows:
                lines.append(" | ".join(cell.text for cell in row.cells))
        text = "\n".join(lines)
        Path(output_path).write_text(text, encoding="utf-8")
        return text
    except Exception as e:
        print(f"  [WARN] DOCX extraction failed: {e}")
        return ""

def synthesize_document(folder_name, transcript, ppt_text, questions_text, output_path):
    """调用 LLM 合成最终培训文档。"""
    import httpx

    # 截断过长内容
    trans_snippet = transcript[:20000] if len(transcript) > 20000 else transcript
    ppt_snippet = ppt_text[:5000] if len(ppt_text) > 5000 else ppt_text

    prompt = f"""请根据以下培训视频转录{f'、PPT讲义' if ppt_text else ''}{f'和培训试题' if questions_text else ''}，直接编写一份完整的培训学习文档。不要加任何开场白或自我介绍。

## 输出结构
1. **培训概述**：主题、设备型号、培训目标
2. **设备整体介绍**：架构、工作流程、核心部件
3. **各模块详细说明**：按培训内容的逻辑顺序展开，保留所有细节和参数
4. **关键参数汇总表**（如有）
{f'5. **试题解答**：逐一详细回答所有试题' if questions_text else ''}
6. **关键结论/FAQ**（如有问答环节）

## 注意事项
- 禁止出现"作为...专家"、"好的，我将..."等AI味开场白，直接输出文档内容
- 不要忽略任何细节和参数说明
- 结合PPT中的术语和框架（如有）
- 试题答案要精准引用视频转录内容
- 不要编造，纯从培训材料中提取
- 中文输出，专业术语保留英文

---

## 视频转录文本

{trans_snippet}

{f'---## PPT讲义内容{ppt_snippet}---' if ppt_text else ''}

{f'---## 培训试题{questions_text}---' if questions_text else ''}
"""
    print(f"  [Synthesize] Prompt: {len(prompt):,} chars, calling LLM...")

    r = httpx.post(
        "https://api.llm.com/v1/chat/completions",
        headers={"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"},
        json={"model": LLM_MODEL, "max_tokens": 8192, "messages": [{"role": "user", "content": prompt}]},
        timeout=httpx.Timeout(600.0, connect=30.0, read=180.0, write=30.0),
    )

    if r.status_code != 200:
        print(f"  [ERROR] LLM API: {r.status_code} {r.text[:200]}")
        return None

    result = r.json()["choices"][0]["message"]["content"]
    Path(output_path).write_text(result, encoding="utf-8")
    return result

def process_folder(folder_path):
    """处理单个培训视频文件夹。"""
    folder_name = os.path.basename(folder_path)
    print(f"\n{'='*60}")
    print(f"Processing: {folder_name}")
    print(f"{'='*60}")

    # 检查是否已处理
    existing_doc = find_file(folder_path, "培训文档.md")
    if existing_doc:
        print(f"  [SKIP] Already has 培训文档.md")
        return True

    # 找文件
    video_file = find_file(folder_path, ".mp4", ".mkv", ".mov", ".avi")
    pptx_file = find_file(folder_path, ".pptx", ".ppt")
    docx_file = find_file(folder_path, ".docx")

    if not video_file:
        print(f"  [ERROR] No video file found")
        return False

    print(f"  Video: {os.path.basename(video_file)} ({os.path.getsize(video_file)/1024/1024:.1f}MB)")
    if pptx_file:
        print(f"  PPT:   {os.path.basename(pptx_file)}")
    if docx_file:
        print(f"  DOCX:  {os.path.basename(docx_file)}")

    stem = Path(video_file).stem

    # Step 0: Extract PPT and DOCX
    ppt_text = ""
    questions_text = ""
    if pptx_file:
        print(f"  [0a] Extracting PPT...")
        ppt_text = extract_pptx(pptx_file, os.path.join(folder_path, f"{stem}_ppt.txt"))
        print(f"       {len(ppt_text)} chars")
    if docx_file:
        print(f"  [0b] Extracting DOCX...")
        questions_text = extract_docx(docx_file, os.path.join(folder_path, f"{stem}_questions.txt"))
        print(f"       {len(questions_text)} chars")

    # Step 1: Extract audio
    print(f"  [1/4] Extracting audio...")
    t0 = time.time()
    audio_path = extract_audio(video_file, output_dir=folder_path)
    duration = get_audio_duration(audio_path)
    print(f"       Done: {duration:.0f}s ({duration/60:.1f}min)")

    # Step 2: Transcribe
    print(f"  [2/4] Transcribing (model={MODEL_SIZE})...")
    transcript = transcribe(audio_path, model_size=MODEL_SIZE, language=LANGUAGE, hf_mirror=HF_MIRROR)
    transcript_file = os.path.join(folder_path, f"{stem}_transcript.txt")
    timed_file = os.path.join(folder_path, f"{stem}_transcript_timed.txt")

    Path(transcript_file).write_text(transcript["text"], encoding="utf-8")
    timed_text = format_transcript_with_timestamps(transcript["segments"])
    Path(timed_file).write_text(timed_text, encoding="utf-8")
    print(f"       {len(transcript['text']):,} chars, {len(transcript['segments'])} segments")

    # Step 3: Synthesize
    print(f"  [3/4] Synthesizing final document...")
    doc_output = os.path.join(folder_path, f"{stem}_培训文档.md")
    result = synthesize_document(
        folder_name,
        transcript["text"],
        ppt_text,
        questions_text,
        doc_output,
    )

    if not result:
        print(f"  [ERROR] Synthesis failed")
        return False

    print(f"       Document: {len(result):,} chars")

    # Cleanup audio
    if os.path.exists(audio_path):
        os.remove(audio_path)
        print(f"  [Clean] Removed temp audio")

    elapsed = time.time() - t0
    print(f"  [DONE] {elapsed:.0f}s total")
    return True

def main():
    base = sys.argv[1] if len(sys.argv) > 1 else r"<training-videos-directory>"
    print(f"Batch processing: {base}")
    print(f"Model: {MODEL_SIZE}, LLM: {LLM_MODEL}")

    dirs = sorted([
        d for d in os.listdir(base)
        if os.path.isdir(os.path.join(base, d))
    ])

    success = 0
    skipped = 0
    failed = []

    for i, d in enumerate(dirs):
        folder = os.path.join(base, d)
        try:
            ok = process_folder(folder)
            if ok and "SKIP" not in str(ok):
                success += 1
            elif ok:
                skipped += 1
            else:
                failed.append(d)
        except Exception as e:
            print(f"  [FATAL] {e}")
            failed.append(d)

    print(f"\n{'='*60}")
    print(f"BATCH COMPLETE: {success} done, {skipped} skipped, {len(failed)} failed")
    if failed:
        print(f"Failed: {failed}")
    print(f"{'='*60}")

if __name__ == "__main__":
    main()
