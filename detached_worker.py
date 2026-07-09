#!/usr/bin/env python3
"""
Detached worker: runs the full pipeline for a list of video folders as a
completely independent process (no parent-child relationship with the launcher).

Usage: python detached_worker.py "folder1" "folder2" ...

Writes progress to {folder}/_worker.log and creates {folder}/_DONE when finished.
"""
import os, sys, time, traceback, subprocess
from pathlib import Path
from datetime import datetime

PROJECT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(PROJECT_DIR))

API_KEY = "YOUR_API_KEY"
LLM_MODEL = "llm-chat"
MODEL_SIZE = "base"

TOPIC_PROMPT = """总结以下培训内容。只关注：运动控制/视觉系统/IO输入输出/传感器/Mapping标定/EAP/软件框架(用户管理/日志管理/参数文件管理)/半导体封装工艺。无关内容跳过。中文输出，保留专业术语英文。

---
{transcript}
{pptx_block}"""


def log(folder, msg):
    ts = datetime.now().strftime("%H:%M:%S")
    line = f"[{ts}] {msg}"
    print(line, flush=True)
    try:
        with open(os.path.join(folder, "_worker.log"), "a", encoding="utf-8") as f:
            f.write(line + "\n")
    except Exception:
        pass


def find_file(folder, *exts):
    for f in sorted(os.listdir(folder)):
        if any(f.lower().endswith(ext) for ext in exts):
            return os.path.join(folder, f)
    return None


def process_one(folder_path):
    folder_name = os.path.basename(folder_path)
    log(folder_path, f"START: {folder_name}")

    # Skip if done
    done_file = os.path.join(folder_path, "_DONE")
    if os.path.exists(done_file):
        log(folder_path, "SKIP: _DONE exists")
        return True

    existing_doc = find_file(folder_path, "_培训文档.md", "培训文档.md")
    if existing_doc:
        # Mark done
        Path(done_file).touch()
        log(folder_path, f"SKIP: doc exists ({os.path.basename(existing_doc)})")
        return True

    video_file = find_file(folder_path, ".mp4", ".mkv", ".mov", ".avi")
    pptx_file = find_file(folder_path, ".pptx")

    transcript_text = None
    stem = None

    if video_file:
        stem = Path(video_file).stem
        log(folder_path, f"Video: {Path(video_file).name}")

        trans_file = find_file(folder_path, "_transcript.txt")
        if trans_file and "_timed" not in trans_file:
            transcript_text = Path(trans_file).read_text(encoding="utf-8")
            log(folder_path, f"Reuse transcript: {len(transcript_text):,} chars")
        else:
            # Transcribe in a subprocess for extra isolation
            transcript_text = _transcribe_in_subprocess(folder_path, video_file, stem)
            if not transcript_text:
                log(folder_path, "FAIL: transcription failed")
                return False
    else:
        log(folder_path, "No video found")
        stem = folder_name

    # Extract PPTX
    pptx_text = ""
    if pptx_file:
        try:
            from pptx import Presentation
            prs = Presentation(pptx_file)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)
            pptx_text = "\n".join(lines)[:8000]
            log(folder_path, f"PPTX: {len(pptx_text)} chars")
        except Exception as e:
            log(folder_path, f"PPTX error: {e}")

    # Synthesize
    if not transcript_text:
        log(folder_path, "No transcript available")
        return True

    import httpx

    trans_snippet = transcript_text[:25000] if len(transcript_text) > 25000 else transcript_text
    pptx_block = f"\n---\nPPT讲义:\n{pptx_text}\n" if pptx_text else ""

    prompt = TOPIC_PROMPT.format(transcript=trans_snippet, pptx_block=pptx_block)
    log(folder_path, f"Synthesizing (prompt {len(prompt):,} chars)...")

    for attempt in range(3):
        try:
            r = httpx.post(
                "https://api.llm.com/v1/chat/completions",
                headers={"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"},
                json={"model": LLM_MODEL, "max_tokens": 8192, "messages": [{"role": "user", "content": prompt}]},
                timeout=httpx.Timeout(600.0),
            )
            if r.status_code == 200:
                result = r.json()["choices"][0]["message"]["content"]
                doc_out = os.path.join(folder_path, f"{stem}_培训文档.md")
                Path(doc_out).write_text(result, encoding="utf-8")
                Path(done_file).touch()
                log(folder_path, f"DONE: {len(result):,} chars")
                return True
            else:
                log(folder_path, f"API error (attempt {attempt+1}): {r.status_code}")
        except Exception as e:
            log(folder_path, f"API exception (attempt {attempt+1}): {e}")
        time.sleep(2)

    log(folder_path, "FAIL: API failed after 3 attempts")
    return False


def _transcribe_in_subprocess(folder_path, video_file, stem):
    """Run transcription in a separate process for memory isolation."""
    import json, tempfile

    # Use a temp JSON file to pass results back
    result_file = os.path.join(folder_path, "_transcribe_result.json")

    transcribe_script = f'''
import sys, os, time, json
sys.path.insert(0, r"{PROJECT_DIR}")
from src.audio import extract_audio, get_audio_duration
from src.transcribe import transcribe, format_transcript_with_timestamps

folder = r"{folder_path}"
video = r"{video_file}"
stem = "{stem}"
model_size = "{MODEL_SIZE}"
result_file = r"{result_file}"

try:
    # Extract audio
    audio = extract_audio(video, output_dir=folder)
    dur = get_audio_duration(audio)
    t0 = time.time()
    t = transcribe(audio, model_size=model_size, language="zh", hf_mirror=True)
    dt = time.time() - t0

    trans_file = os.path.join(folder, f"{{stem}}_transcript.txt")
    timed_file = os.path.join(folder, f"{{stem}}_transcript_timed.txt")
    with open(trans_file, "w", encoding="utf-8") as f:
        f.write(t["text"])
    with open(timed_file, "w", encoding="utf-8") as f:
        f.write(format_transcript_with_timestamps(t["segments"]))

    if os.path.exists(audio):
        os.remove(audio)

    with open(result_file, "w", encoding="utf-8") as f:
        json.dump({{"status": "ok", "chars": len(t["text"]), "elapsed": dt}}, f)
    sys.exit(0)
except Exception as e:
    with open(result_file, "w", encoding="utf-8") as f:
        json.dump({{"status": "error", "message": str(e)}}, f)
    sys.exit(1)
'''

    log(folder_path, "Launching transcription subprocess...")
    proc = subprocess.Popen(
        [sys.executable, "-c", transcribe_script],
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
        creationflags=subprocess.CREATE_NO_WINDOW if sys.platform == "win32" else 0,
    )

    # Wait for subprocess to complete (it writes result_file when done)
    log(folder_path, f"Waiting for transcription (pid={proc.pid})...")
    proc.wait(timeout=7200)  # 2 hour max

    if os.path.exists(result_file):
        result = json.loads(Path(result_file).read_text(encoding="utf-8"))
        os.remove(result_file)
        if result["status"] == "ok":
            # Read transcript back
            trans_file = os.path.join(folder_path, f"{stem}_transcript.txt")
            text = Path(trans_file).read_text(encoding="utf-8")
            log(folder_path, f"Transcribed: {len(text):,} chars, {result['elapsed']:.0f}s")
            return text

    log(folder_path, "Transcription subprocess failed")
    return None


def main():
    folders = sys.argv[1:]
    if not folders:
        print("Usage: python detached_worker.py <folder1> [folder2 ...]")
        sys.exit(1)

    # Write PID file so we know it's running
    pid_file = PROJECT_DIR / "_worker.pid"
    pid_file.write_text(str(os.getpid()))

    total = len(folders)
    success = 0
    for i, folder in enumerate(folders):
        if not os.path.isdir(folder):
            print(f"[{i+1}/{total}] MISSING: {folder}", flush=True)
            continue
        try:
            ok = process_one(folder)
            if ok:
                success += 1
        except Exception:
            traceback.print_exc()

    # Cleanup
    if pid_file.exists():
        pid_file.unlink()

    print(f"\nALL DONE: {success}/{total} succeeded", flush=True)


if __name__ == "__main__":
    main()
