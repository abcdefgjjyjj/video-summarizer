#!/usr/bin/env python3
"""
Robust batch v2: Process remaining videos via multiprocessing transcription (2 workers),
then threaded synthesis (6 workers). Designed for resilience.

Usage: python robust_batch.py
"""
import os, sys, time
from pathlib import Path
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor, as_completed

PROJECT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(PROJECT_DIR))

API_KEY = "YOUR_API_KEY"
MODEL_SIZE = "base"
LLM_MODEL = "llm-chat"
HF_MIRROR = True
LANGUAGE = "zh"

RELEVANT = [
    "31.EAP介绍和调试过程讲解-员工E", "33.新版视觉2.0软件操作简介-员工B",
    "34.建模注意事项-员工H", "37.1210时序规划与分析-员工J",
    "39.基于1210的通讯介绍-员工K", "40.AI技术在芯片检测中应用-员工B",
    "41.1250设备常见问题及其分析方法（一）-员工A", "42.1250设备常见问题及其分析方法（二）-员工A",
    "43.1250设备常见问题及其分析方法（三）-员工A", "44.设备模块的封装和应用-员工B",
    "48.1220框架讲解-员工C", "49.BGA封装流程介绍-员工D",
    "52.EAP的软件安装过程和四大模块的应用-员工E", "53.视觉2新功能介绍和标定方式-员工B",
    "54.Memory工艺培训", "56.EAP模拟器和EAP代码的应用-员工E",
    "57. 视觉基础知识 及 镜头相机选型 - 员工F", "59.MIT视觉GP软件培训-Colleague G",
    "60.视觉GP软件培训-员工H", "61.FC1250时序分析-员工I",
]

TOPIC_PROMPT = """总结以下培训内容。只关注：运动控制(Motion)/视觉(Vision)/IO/传感器/Mapping标定/EAP/EFEM/用户管理/日志管理/参数文件管理/半导体封装工艺。无关内容跳过。中文输出，保留专业术语英文。

---
{transcript}
{pptx_block}
"""


def _find_video(folder):
    for f in sorted(os.listdir(folder)):
        if any(f.lower().endswith(ext) for ext in ('.mp4', '.mkv', '.mov', '.avi')):
            return os.path.join(folder, f)
    return None


def transcribe_worker(folder_path):
    """CPU-bound: extract audio + transcribe"""
    import sys
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from src.audio import extract_audio, get_audio_duration
    from src.transcribe import transcribe, format_transcript_with_timestamps

    folder_name = os.path.basename(folder_path)
    video_file = _find_video(folder_path)
    if not video_file:
        return {"folder": folder_name, "status": "no_video"}

    stem = Path(video_file).stem
    doc_file = os.path.join(folder_path, f"{stem}_培训文档.md")
    if os.path.exists(doc_file):
        return {"folder": folder_name, "status": "skipped"}

    t0 = time.time()
    try:
        audio_path = extract_audio(video_file, output_dir=folder_path)
        duration = get_audio_duration(audio_path)
        transcript = transcribe(audio_path, model_size=MODEL_SIZE, language=LANGUAGE, hf_mirror=HF_MIRROR)
        trans_file = os.path.join(folder_path, f"{stem}_transcript.txt")
        timed_file = os.path.join(folder_path, f"{stem}_transcript_timed.txt")
        Path(trans_file).write_text(transcript["text"], encoding="utf-8")
        Path(timed_file).write_text(format_transcript_with_timestamps(transcript["segments"]), encoding="utf-8")
        if os.path.exists(audio_path):
            os.remove(audio_path)
        elapsed = time.time() - t0
        return {
            "folder": folder_name, "status": "transcribed",
            "chars": len(transcript["text"]), "duration_min": round(duration/60, 1),
            "elapsed_s": round(elapsed, 0), "stem": stem,
        }
    except Exception as e:
        return {"folder": folder_name, "error": str(e)}


def synthesize_worker(args):
    """I/O-bound: call LLM API"""
    folder_path, api_key, model = args
    import httpx

    folder_name = os.path.basename(folder_path)

    # Find transcript
    transcript_text = ""
    stem = None
    for f in sorted(os.listdir(folder_path)):
        if f.endswith("_transcript.txt") and not f.endswith("_timed.txt"):
            transcript_text = Path(os.path.join(folder_path, f)).read_text(encoding="utf-8")
            stem = Path(f).stem.replace("_transcript", "")
            break

    if not transcript_text:
        return {"folder": folder_name, "error": "No transcript found"}

    # Extract PPTX
    pptx_text = ""
    for f in sorted(os.listdir(folder_path)):
        if f.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(os.path.join(folder_path, f))
                lines = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = para.text.strip()
                                if t:
                                    lines.append(t)
                pptx_text = "\n".join(lines)[:8000]
            except Exception:
                pass
            break

    trans_snippet = transcript_text[:25000] if len(transcript_text) > 25000 else transcript_text
    pptx_block = f"\n---\nPPT讲义:\n{pptx_text}\n" if pptx_text else ""

    prompt = TOPIC_PROMPT.format(transcript=trans_snippet, pptx_block=pptx_block)

    try:
        r = httpx.post(
            "https://api.llm.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": model, "max_tokens": 8192, "messages": [{"role": "user", "content": prompt}]},
            timeout=httpx.Timeout(600.0, connect=30.0, read=300.0, write=30.0),
        )
        if r.status_code != 200:
            return {"folder": folder_name, "error": f"API {r.status_code}"}
        result = r.json()["choices"][0]["message"]["content"]
        doc_output = os.path.join(folder_path, f"{stem}_培训文档.md")
        Path(doc_output).write_text(result, encoding="utf-8")
        return {"folder": folder_name, "status": "done", "chars": len(result)}
    except Exception as e:
        return {"folder": folder_name, "error": str(e)}


def main():
    BASE = r"<training-videos-directory>"
    folders = [os.path.join(BASE, t) for t in RELEVANT if os.path.isdir(os.path.join(BASE, t))]

    # Filter out already done
    unprocessed = []
    for fp in folders:
        has_doc = any(f.endswith("_培训文档.md") or f.endswith("培训文档.md") for f in os.listdir(fp))
        if not has_doc:
            unprocessed.append(fp)

    if not unprocessed:
        print("All done!")
        return

    print(f"{'='*60}")
    print(f"Robust Batch v2: {len(unprocessed)} videos remaining")
    print(f"{'='*60}")
    for f in unprocessed:
        print(f"  {os.path.basename(f)}")

    # Phase 1: Parallel transcription (2 workers)
    print(f"\nPhase 1: Parallel transcription (2 workers)")
    print(f"{'='*60}")
    t0 = time.time()
    transcribed = []

    with ProcessPoolExecutor(max_workers=2) as executor:
        futures = {executor.submit(transcribe_worker, f): f for f in unprocessed}
        for future in as_completed(futures):
            r = future.result()
            if r.get("status") == "transcribed":
                print(f"  [{r['folder']}] OK {r['duration_min']}min {r['chars']:,}chars {r['elapsed_s']}s")
                transcribed.append(r["folder"])
            elif r.get("status") == "skipped":
                print(f"  [{r['folder']}] SKIP")
                transcribed.append(r["folder"])
            else:
                print(f"  [{r['folder']}] FAIL: {r.get('error', '?')}")
                transcribed.append(r["folder"])  # Try synthesis anyway

    t1 = time.time()
    print(f"Phase 1 done: {t1-t0:.0f}s")

    # Phase 2: Parallel synthesis (6 threads)
    print(f"\nPhase 2: Parallel synthesis (6 threads)")
    print(f"{'='*60}")

    tasks = [(os.path.join(BASE, fn) if not fn.startswith(BASE) else fn, API_KEY, LLM_MODEL) for fn in transcribed]
    with ThreadPoolExecutor(max_workers=6) as executor:
        futures = {executor.submit(synthesize_worker, t): t for t in tasks}
        for future in as_completed(futures):
            r = future.result()
            if r.get("status") == "done":
                print(f"  [{r['folder']}] DONE {r['chars']:,} chars")
            else:
                print(f"  [{r['folder']}] FAIL: {r.get('error', '?')}")

    t2 = time.time()
    print(f"\nPhase 2 done: {t2-t1:.0f}s")
    print(f"Total: {t2-t0:.0f}s")


if __name__ == "__main__":
    main()
