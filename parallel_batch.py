#!/usr/bin/env python3
"""
Parallel batch — Phase 1: 2-process transcription, Phase 2: 6-thread synthesis.
"""
import os, sys, time, json
from pathlib import Path
from concurrent.futures import ProcessPoolExecutor, ThreadPoolExecutor, as_completed

PROJECT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(PROJECT_DIR))

API_KEY = "YOUR_API_KEY"
MODEL_SIZE = "base"
LLM_MODEL = "llm-chat"
HF_MIRROR = True
LANGUAGE = "zh"


# ---- worker functions (must be module-level for Windows multiprocessing) ----

def _find_video(folder_path: str) -> str | None:
    for f in sorted(os.listdir(folder_path)):
        if any(f.lower().endswith(ext) for ext in ('.mp4', '.mkv', '.mov', '.avi')):
            return os.path.join(folder_path, f)
    return None


def transcribe_worker(folder_path: str) -> dict:
    """CPU-bound worker — extract audio + transcribe. Runs in subprocess."""
    import sys
    from pathlib import Path
    sys.path.insert(0, str(Path(folder_path).parent.parent / "video-summarizer"))

    from src.audio import extract_audio, get_audio_duration
    from src.transcribe import transcribe, format_transcript_with_timestamps

    folder_name = os.path.basename(folder_path)
    video_file = _find_video(folder_path)
    if not video_file:
        return {"folder": folder_name, "error": "No video found"}

    stem = Path(video_file).stem
    doc_file = os.path.join(folder_path, f"{stem}_培训文档.md")
    if os.path.exists(doc_file):
        return {"folder": folder_name, "status": "skipped", "reason": "doc exists"}

    t0 = time.time()
    try:
        audio_path = extract_audio(video_file, output_dir=folder_path)
        duration = get_audio_duration(audio_path)
        transcript = transcribe(audio_path, model_size=MODEL_SIZE, language=LANGUAGE, hf_mirror=HF_MIRROR)

        # Save transcripts
        transcript_file = os.path.join(folder_path, f"{stem}_transcript.txt")
        timed_file = os.path.join(folder_path, f"{stem}_transcript_timed.txt")
        Path(transcript_file).write_text(transcript["text"], encoding="utf-8")
        Path(timed_file).write_text(format_transcript_with_timestamps(transcript["segments"]), encoding="utf-8")

        # Cleanup audio
        if os.path.exists(audio_path):
            os.remove(audio_path)

        elapsed = time.time() - t0
        return {
            "folder": folder_name,
            "status": "transcribed",
            "chars": len(transcript["text"]),
            "duration_min": round(duration / 60, 1),
            "elapsed_s": round(elapsed, 0),
            "stem": stem,
        }
    except Exception as e:
        return {"folder": folder_name, "error": str(e)}


def synthesize_worker(args: tuple) -> dict:
    """I/O-bound worker — call LLM API. Runs in thread."""
    folder_path, api_key, model = args
    import httpx

    folder_name = os.path.basename(folder_path)

    # Find transcript and PPT/DOCX
    transcript_file = None
    ppt_text = ""
    questions_text = ""
    for f in sorted(os.listdir(folder_path)):
        if f.endswith("_transcript.txt") and not f.endswith("_timed.txt"):
            transcript_file = os.path.join(folder_path, f)
        elif f.endswith(".txt") and "ppt" in f.lower():
            ppt_text = Path(os.path.join(folder_path, f)).read_text(encoding="utf-8")[:5000]
        elif f.endswith(".txt") and "question" in f.lower():
            questions_text = Path(os.path.join(folder_path, f)).read_text(encoding="utf-8")

    if not transcript_file:
        # Check for PPTX/DOCX
        for f in sorted(os.listdir(folder_path)):
            fp = os.path.join(folder_path, f)
            if f.endswith(".pptx"):
                try:
                    from pptx import Presentation
                    prs = Presentation(fp)
                    ppt_text = "\n".join(
                        p.text.strip() for slide in prs.slides
                        for shape in slide.shapes if shape.has_text_frame
                        for p in shape.text_frame.paragraphs if p.text.strip()
                    )[:5000]
                except Exception:
                    pass
            elif f.endswith(".docx"):
                try:
                    from docx import Document
                    doc = Document(fp)
                    parts = []
                    for para in doc.paragraphs:
                        if para.text.strip():
                            parts.append(para.text.strip())
                    for table in doc.tables:
                        for row in table.rows:
                            parts.append(" | ".join(cell.text for cell in row.cells))
                    questions_text = "\n".join(parts)
                except Exception:
                    pass

    stem = Path(transcript_file).stem.replace("_transcript", "")
    transcript_text = Path(transcript_file).read_text(encoding="utf-8")
    trans_snippet = transcript_text[:25000] if len(transcript_text) > 25000 else transcript_text

    prompt = f"""请根据以下培训视频转录{f'、PPT讲义' if ppt_text else ''}{f'和培训试题' if questions_text else ''}，直接编写一份完整的培训学习文档。不要加任何开场白或自我介绍。

## 输出结构
1. **培训概述**：主题、设备型号、培训目标
2. **设备整体介绍**：架构、工作流程、核心部件
3. **各模块详细说明**：按培训内容的逻辑顺序展开，保留所有细节和参数
4. **关键参数汇总表**（如有）
{f'5. **试题解答**：逐一详细回答所有试题' if questions_text else ''}
6. **关键结论/FAQ**（如有问答环节）

## 注意事项
- 禁止出现"作为...专家"、"好的，我将..."等AI味开场白，直接输出文档内容
- 不要忽略任何细节和参数说明
- 结合PPT中的术语和框架（如有）
- 试题答案要精准引用视频转录内容
- 不要编造，纯从培训材料中提取
- 中文输出，专业术语保留英文

---

## 视频转录文本

{trans_snippet}

{f'---## PPT讲义内容{ppt_text}---' if ppt_text else ''}

{f'---## 培训试题{questions_text}---' if questions_text else ''}
"""
    try:
        r = httpx.post(
            "https://api.llm.com/v1/chat/completions",
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json={"model": model, "max_tokens": 8192, "messages": [{"role": "user", "content": prompt}]},
            timeout=httpx.Timeout(600.0, connect=30.0, read=300.0, write=30.0),
        )
        if r.status_code != 200:
            return {"folder": folder_name, "error": f"API {r.status_code}: {r.text[:100]}"}

        result = r.json()["choices"][0]["message"]["content"]
        doc_output = os.path.join(folder_path, f"{stem}_培训文档.md")
        Path(doc_output).write_text(result, encoding="utf-8")
        return {"folder": folder_name, "status": "done", "chars": len(result)}
    except Exception as e:
        return {"folder": folder_name, "error": str(e)}


# ---- main ----

def main():
    base = sys.argv[1] if len(sys.argv) > 1 else r"<training-videos-directory>"

    # Discover unprocessed folders
    dirs = sorted([
        os.path.join(base, d) for d in os.listdir(base)
        if os.path.isdir(os.path.join(base, d))
    ])

    unprocessed = []
    for folder_path in dirs:
        has_doc = any(
            f.endswith("_培训文档.md") or f.endswith("培训文档.md")
            for f in os.listdir(folder_path)
        )
        if not has_doc:
            unprocessed.append(folder_path)

    if not unprocessed:
        print("All folders already processed!")
        return

    print(f"Unprocessed folders ({len(unprocessed)}):")
    for f in unprocessed:
        print(f"  - {os.path.basename(f)}")

    # ---- Phase 1: Parallel transcription (2 workers, CPU-bound) ----
    print(f"\n{'='*60}")
    print("Phase 1: Parallel transcription (2 workers)")
    print(f"{'='*60}")
    t0 = time.time()

    with ProcessPoolExecutor(max_workers=2) as executor:
        futures = {executor.submit(transcribe_worker, f): f for f in unprocessed}
        for future in as_completed(futures):
            r = future.result()
            status = r.get("status") or r.get("error", "unknown")
            if r.get("status") == "transcribed":
                print(f"  [{r['folder']}] OK  {r['duration_min']}min audio, {r['chars']:,} chars, {r['elapsed_s']:.0f}s")
            elif r.get("status") == "skipped":
                print(f"  [{r['folder']}] SKIP ({r.get('reason')})")
            else:
                print(f"  [{r['folder']}] FAIL: {status}")

    t1 = time.time()
    print(f"Phase 1 done: {t1 - t0:.0f}s")

    # ---- Phase 2: Parallel synthesis (6 workers, I/O-bound) ----
    print(f"\n{'='*60}")
    print("Phase 2: Parallel synthesis (6 threads)")
    print(f"{'='*60}")

    with ThreadPoolExecutor(max_workers=6) as executor:
        futures = {
            executor.submit(synthesize_worker, (f, API_KEY, LLM_MODEL)): f
            for f in unprocessed
        }
        for future in as_completed(futures):
            r = future.result()
            if r.get("status") == "done":
                print(f"  [{r['folder']}] DONE  {r['chars']:,} chars")
            else:
                print(f"  [{r['folder']}] FAIL: {r.get('error', 'unknown')}")

    t2 = time.time()
    print(f"\nPhase 2 done: {t2 - t1:.0f}s")
    print(f"Total: {t2 - t0:.0f}s")


if __name__ == "__main__":
    main()
