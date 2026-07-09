#!/usr/bin/env python3
"""Process a single video folder: transcribe + synthesize. Standalone subprocess for memory isolation."""
import os, sys, time, gc, traceback
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

API_KEY = "YOUR_API_KEY"
LLM_MODEL = "llm-chat"
MODEL_SIZE = "base"

TOPIC_PROMPT = """总结以下培训内容。只关注：运动控制/视觉系统/IO输入输出/传感器/Mapping标定/EAP/软件框架(用户管理/日志管理/参数文件管理)/半导体封装工艺。无关内容跳过。中文输出，保留专业术语英文。

---
{transcript}
{pptx_block}"""


def find_file(folder, *exts):
    for f in sorted(os.listdir(folder)):
        if any(f.lower().endswith(ext) for ext in exts):
            return os.path.join(folder, f)
    return None


def process_one(folder_path):
    folder_name = os.path.basename(folder_path)
    print(f"\n=== {folder_name} ===", flush=True)

    # Skip if done
    done = find_file(folder_path, "_培训文档.md", "培训文档.md")
    if done:
        print(f"  SKIP: already has doc", flush=True)
        return True

    video_file = find_file(folder_path, ".mp4", ".mkv", ".mov", ".avi")
    pptx_file = find_file(folder_path, ".pptx")

    transcript_text = None
    stem = None

    if video_file:
        stem = Path(video_file).stem
        print(f"  Video: {Path(video_file).name}", flush=True)

        # Check for existing transcript
        trans_file = find_file(folder_path, "_transcript.txt")
        if trans_file and "_timed" not in trans_file:
            transcript_text = Path(trans_file).read_text(encoding="utf-8")
            print(f"  Reuse transcript: {len(transcript_text):,} chars", flush=True)
        else:
            # Transcribe
            from src.audio import extract_audio, get_audio_duration
            from src.transcribe import transcribe, format_transcript_with_timestamps

            print(f"  Extracting audio...", flush=True)
            audio = extract_audio(video_file, output_dir=folder_path)
            dur = get_audio_duration(audio)
            print(f"  Audio: {dur:.0f}s, transcribing...", flush=True)

            t0 = time.time()
            t = transcribe(audio, model_size=MODEL_SIZE, language="zh", hf_mirror=True)
            dt = time.time() - t0

            trans_file = os.path.join(folder_path, f"{stem}_transcript.txt")
            timed_file = os.path.join(folder_path, f"{stem}_transcript_timed.txt")
            Path(trans_file).write_text(t["text"], encoding="utf-8")
            Path(timed_file).write_text(format_transcript_with_timestamps(t["segments"]), encoding="utf-8")
            transcript_text = t["text"]
            print(f"  Transcript: {len(transcript_text):,} chars, {dt:.0f}s", flush=True)

            if os.path.exists(audio):
                os.remove(audio)
            # Force cleanup
            del t, audio
            gc.collect()
    else:
        print(f"  No video found", flush=True)
        stem = folder_name

    # Extract PPTX
    pptx_text = ""
    if pptx_file:
        try:
            from pptx import Presentation
            prs = Presentation(pptx_file)
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                lines.append(t)
            pptx_text = "\n".join(lines)[:8000]
            print(f"  PPTX: {len(pptx_text)} chars", flush=True)
        except Exception as e:
            print(f"  PPTX error: {e}", flush=True)

    # Synthesize
    if not transcript_text:
        print(f"  No transcript available", flush=True)
        return True

    import httpx

    trans_snippet = transcript_text[:25000] if len(transcript_text) > 25000 else transcript_text
    pptx_block = f"\n---\nPPT讲义:\n{pptx_text}\n" if pptx_text else ""

    prompt = TOPIC_PROMPT.format(transcript=trans_snippet, pptx_block=pptx_block)
    print(f"  Synthesizing (prompt {len(prompt):,} chars)...", flush=True)

    r = httpx.post(
        "https://api.llm.com/v1/chat/completions",
        headers={"Authorization": f"Bearer {API_KEY}", "Content-Type": "application/json"},
        json={"model": LLM_MODEL, "max_tokens": 8192, "messages": [{"role": "user", "content": prompt}]},
        timeout=httpx.Timeout(600.0),
    )

    if r.status_code == 200:
        result = r.json()["choices"][0]["message"]["content"]
        doc_out = os.path.join(folder_path, f"{stem}_培训文档.md")
        Path(doc_out).write_text(result, encoding="utf-8")
        print(f"  DONE: {len(result):,} chars", flush=True)
        return True
    else:
        print(f"  API ERROR: {r.status_code} {r.text[:200]}", flush=True)
        return False


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python process_one.py <folder_path>")
        sys.exit(1)

    folder = sys.argv[1]
    try:
        success = process_one(folder)
        sys.exit(0 if success else 1)
    except Exception as e:
        traceback.print_exc()
        sys.exit(1)
